#!/usr/bin/env python
# -*- coding: utf-8 -*-

import numpy as np
import os, functools
import pptx 

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR 
from pptx.dml.color import RGBColor
import uuid

import PIModelData
import Plot as piplot



class PIController(object):
    def __init__(self, model, view):
        self.model = model
        self.view = view

    def save(self, filename):
        v = self.view
        m = self.model

        v.change_title(m.title_data())
        v.change_HVConsumption(m.hvConsumption_data())

        for data in m.plot_data():
            v.add_pi_plot(data['plot_data'], data['position'])

        for data in m.ap_data():
            v.add_apbox(data['value'], data['position'])

        v.save(filename)

class PISlide(object):
    """
    View Object for powerpoint
    """
    def __init__(self):
        powerpoint_template = os.path.join(os.path.dirname(__file__), "base_pi.pptx")
        self.presentation = pptx.Presentation(powerpoint_template)
    
    def _create_text_box(self):
        pass


    def save(self, filename):
        dirname = os.path.dirname(filename)

        self.presentation.save(filename)

    def add_pi_plot(self, plot_data, position):
        slide = self.presentation.slides[0]

        unique_filename = unicode(uuid.uuid4()) + u".png"
        new_fig = piplot.plot_pi(plot_data)
        new_fig.save(unique_filename)

        pic = slide.shapes.add_picture(unique_filename, position["left"], 
                                      position["top"],
                                      position["width"], 
                                      position["height"], )
        try:
            os.remove(unique_filename)
        except Exception as err:
            print err

        return pic 


    def add_apbox(self, value, position):
        base_string = (u"Arbeitspunkt:  t = {0} s;   n = {1} 1/min;   RCI "
                          u"= {2} %;   T = {3} °C;   SoH = {4} %;  {5}") 
        result_string = base_string.format(value['t'], value['n'], value['RCI'], 
                                           value['T'], value['SoH'], 
                                           value['mode'], )
        slide = self.presentation.slides[0]

        shape = slide.shapes.add_textbox(position["left"], 
                                         position["top"],
                                         position["width"], 
                                         position["height"], )
        shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        shape.fill.background()

        tf = shape.text_frame 
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.clear()

        p = tf.paragraphs[0]

        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = result_string
        run.font.name = 'BMW Group Condensed'
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run.font.bold = False

        return shape

    def change_title(self, value):
        shape_num = self._match_shape_number(u"PI_MTITLE")
        self._change_text_at_shape(slide_num = 0, shape_num = shape_num, text = value)

    def change_HVConsumption(self, value):
        number_float = float(value)
        formated_string = "{:.1f}".format(number_float) 
        resulted_string = unicode(formated_string).replace(".",",") 
        shape_num = self._match_shape_number(u"PI_STITLE")
        subtitle_text = u"(inkl. " + resulted_string +  u" kW Bordnetzbedarf)"
        self._change_text_at_shape(slide_num = 0, shape_num = shape_num, text = subtitle_text)

    def _match_shape_number(self, text, slide_num = 0):
        all_matching_number = [i for i, shape in enumerate(self.presentation.slides[slide_num].shapes) if self._read_text(shape) == unicode(text)]
        return all_matching_number[0] 

    @staticmethod
    def _read_text(shape):
        if shape.has_text_frame:
            return unicode(shape.text_frame.text)
        else:
            return u""

    def _change_text_at_shape(self, slide_num, shape_num, text):
        self.presentation.slides[slide_num].shapes[shape_num].text_frame.paragraphs[0].runs[0].text = text
            
class PIDummyMasterModel(object):
    def __init__(self):
        pass

    def plot_data(self):
        positions = PIModelData.get_plot_positions()
        line_data = dict()
        x1 = np.arange(0.0, 1.1, 0.1) * 1.0324
        y = np.arange(0.0, 22, 2)* 1.0324 + 123
        line_data['x'] = x1
        line_data['x_label'] = u'RCI (1)'
        line_data['y'] = y
        line_data['y_label'] = '$P_{mech}$ (kW)'
        line_data['y_unit'] = 'kW'
        line_data['color'] = np.arange(0.0, 11, 1)
        text = ['160kW@30s\n-312A@423V@103A(AC)', '160kW@30s\n-312A@423V@103A(AC)', '160kW@30s\n-312A@423V@103A(AC)']
        annotations = zip(x1[1:10:4], y[1:10:4], text) 

        plot_data = dict(line_data = line_data, 
                         star_annotations = annotations, 
                         y_max = None)

        plot_data_list = [plot_data for _ in range(0,5)]
        return [dict(plot_data=plot_data, position=position) 
                for plot_data, position in zip(plot_data_list, positions)]

    def title_data(self):
        return "Dummy title"

    def hvConsumption_data(self):
        return 1.532

    def ap_data(self):
        positions = PIModelData.get_master_aps_positions()
        value_one = dict(t = 300, n = 2000,  RCI = 0.3, T = -25, SoH = 1, mode = "mot")
        values = [value_one for _ in range(0,3)]
        return [dict(value=value, position=position) for value, position in zip(values, positions)]

class PIDummyModel(object):
    def __init__(self):
        self.pidummy = PIDummyMasterModel()

    def plot_data(self):
        return self.pidummy.plot_data()

    def title_data(self):
        return self.pidummy.title_data()

    def hvConsumption_data(self):
        return self.pidummy.hvConsumption_data()

    def ap_data(self):
        position = PIModelData.get_normal_aps_positions()
        value_one = dict(t = 300, n = 2000,  RCI = 0.3, T = -25, SoH = 1, mode = "mot")
        return [dict(value=value_one, position=position)]

if __name__ == "__main__":
    piview = PISlide()
    pimodel = PIDummyModel()
    picontroller = PIController(pimodel, piview)

    picontroller.save("test_new.pptx")
